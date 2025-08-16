from pathlib import Path
from PIL import Image
import docx
import PyPDF2
from PyPDF2.errors import PdfReadError
import pandas as pd
import imagehash
import pptx
from bs4 import BeautifulSoup
from simhash import Simhash
from difflib import SequenceMatcher
from zipfile import is_zipfile
from concurrent.futures import ThreadPoolExecutor, TimeoutError
from tqdm import tqdm

from file_data.file_info import FileInfo
from file_data.file_type_enum import FileType
from sorter import Sorter
from logger import Logger
from console_writer import ConsoleWriter
from decorators import suppress_stderr

class Hasher():
    def __init__(self, sorter : Sorter, _logger : Logger):
        self._sorter = sorter
        self.logger = _logger

    def count_hashes(self, sorted_files : dict) -> None:
        """Count percentual hash of all files in inpued dict, 
        set hash in each FileData object."""        
        hash_functions = {
            FileType.TEXT : self.extract_text_file,
            FileType.DOCX : self.extract_docx_file,
            FileType.DOC : self.extract_doc_file,
            FileType.PDF : self.extract_pdf_file,
            FileType.SPREADSHEET : self.extract_spreadsheet_file,
            FileType.PRESENTATION : self.extract_ppt_file,
            FileType.HTML : self.extract_simhash_from_htm,
            FileType.IMAGE : self.get_image_phash,
            FileType.OTHER : self.manage_other_hash,
        }

        hashing_info = ConsoleWriter.get_hash_counting_info()

        hashing_info['start']()
        for file_type, one_type_files in sorted_files.items():
            for file_info in tqdm(one_type_files, desc=hashing_info['hashing'](file_type.name), unit="file"):
                file_info.set_hash(hash_functions[file_type](file_info.get_path()))
            one_type_files = self._sorter.sort_by_hash(one_type_files)

        hashing_info['end']()
        return sorted_files

    def _get_simhash_from_text(text : str) -> int:
        """Count simhash from text."""
        if not text:
            return None
        return Simhash(text.split()).value

    def extract_text_file(self, path: Path) -> int:
        """Extract text file, return simhash of a file."""
        text = None
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
        return Hasher._get_simhash_from_text(text)

    def extract_docx_file(self, path: Path) -> str:
        """Extract docx file, return simhash of a file."""
        text = None
        try:
            if Hasher.is_probably_valid_docx(path):
                doc = docx.Document(path)
                text =  "\n".join([para.text for para in doc.paragraphs])
        except Exception as e:
            self.logger.add_to_corrupted(path)
        return Hasher._get_simhash_from_text(text)

    def is_probably_valid_docx(path: Path) -> bool:
        """Check, if file is valid docx."""
        return path.suffix.lower() == ".docx" and is_zipfile(path)
        
    def extract_doc_file(self, path: Path) -> str:
        """Extract content of doc file, count and return simhash."""
        text = None
        pass
        # try:
        #     word = win32com.client.Dispatch("Word.Application")
        #     word.Visible = False
        #     doc = word.Documents.Open(str(path))
        #     try:
        #         text = doc.Content.Text
        #     finally:
        #         doc.Close(False)
        #         word.Quit()

        #     cleaned_text = text.strip().replace("\r", "").replace("\n", "")
        #     return hashlib.sha256(cleaned_text.encode('utf-8')).hexdigest()
        
        # except Exception as e:
        #     self.logger.add_to_corrupted(path)
        #     return ""
        return Hasher._get_simhash_from_text(text)

    def extract_pdf_file(self, path: Path) -> str:
        """Extract pdf file if it is possible and return content as a string."""
        with suppress_stderr():
            try:
                with ThreadPoolExecutor(max_workers=1) as executor:
                    future = executor.submit(Hasher._read_pdf, path)
                    return future.result(timeout=4)    
            except TimeoutError:
                self.logger.add_to_corrupted(path)
                return None
            except PdfReadError as e:
                self.logger.add_to_corrupted(path)
                return None
            except Exception as e:
                self.logger.add_to_corrupted(path)
                return None

    def _read_pdf(path: Path) -> str:
        """Read pdf file, return content as a string."""
        text = ""
        with open(path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text += page.extract_text() or ""
        return text

    def extract_spreadsheet_file(self, path: Path) -> str:
        """Extract spreadsheet file, return content as a string."""
        try:
            if path.suffix.lower() == ".csv":
                df = pd.read_csv(path)
            else:
                df = pd.read_excel(path, engine='openpyxl')
            return df.to_csv(index=False)
        except:
            self.logger.add_to_corrupted(path)
            return None
        
    def extract_ppt_file(self, path: Path) -> str:
        """Extract ppt file (presentation), return content as a string."""
        try:
            prs = pptx.Presentation(path)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return "\n".join(text_runs)
        except Exception as e:
            self.logger.add_to_corrupted(path)
            return f"ERROR: {e}"

    def extract_simhash_from_htm(self, path: Path) -> int:
        """Extract html code, count and return simhash of a file."""
        soup = None 
        try:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                soup = BeautifulSoup(f, "html.parser")
            soup.get_text()
        except Exception as e:
            self.logger.add_to_corrupted(path)
        return Hasher._get_simhash_from_text(soup)

    def get_image_phash(self, path: Path) -> int:
        """count percentual hash of image."""
        try:
            hash = imagehash.average_hash(Image.open(path))
            return hash
        except Exception as e:
            self.logger.add_to_corrupted(path)
            return None

    def manage_other_hash(self, path : Path) -> None:
        """Manage hash computing for FileType.OTHER files."""
        return None

    def _sequence_matcher(hash1: str, hash2: str) -> float:
        """Count percentual similarity of two files."""
        if hash1 == hash2:
            return 1
        return 0

    def _hamming_similarity_simhash(simhash1 : int, simhash2 : int) -> float:
        """Return percentage similarity of two simhashes."""
        hamming_distance = bin(simhash1 ^ simhash2).count('1')
        similarity = (1 - hamming_distance / 64)
        return round(similarity, 2)

    def _hamming_distance_images(hash1 : int, hash2 : int) -> int:
        """Return percentage similarity of two image hashes."""
        hamming_distance = abs(hash1 - hash2)
        similarity = (1 - hamming_distance / 64)
        return round(similarity, 2)

    def similarity_score(hash1, hash2, file_type) -> float:
        """Return similarity score (percentage) of two hashes."""
        if not hash1 or not hash2:
            return 0
        if file_type in ['text', 'doc', 'docx']:
            return Hasher._hamming_similarity_simhash(hash1, hash2)
        elif file_type == 'image':
            return Hasher._hamming_distance_images(hash1, hash2)
        else:
            return Hasher._sequence_matcher(str(hash1), str(hash2))
