from pptx import Presentation
import zipfile
import io
from PIL import Image

from file_data.file_info import FileInfo
from hashers.text_hashers.text_hasher_base import TextHasherBase
from hashers.phash_hashers.phash_hasher_base import PhashHasherBase
from decorators import suppress_stderr
from config import PDF_IMAGE_MAX_PAGES

class PresentationHasher(TextHasherBase, PhashHasherBase):
    def __init__(self, sorter, logger):
        super().__init__(sorter, logger)

    def extract_hashes(self, file_info : FileInfo):
        """Extract text and image hashes of the file, set hashes into file_info object."""
        text = self._extract_text(file_info)
        file_info.set_text(text)
        file_info.set_text_hash(self._extract_hash_from_text(text))
        file_info.set_image_hash(self._extract_image_hash(file_info))

    def _extract_text(self, file_info : FileInfo) -> str:
        """Extract ppt file (presentation), return content as a string."""
        self._handle_unsupported_ppt_exception(file_info)
        with suppress_stderr():
            try:
                prs = Presentation(file_info.get_path())
                all_texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text') and shape.text:
                            all_texts.append(shape.text)
                combined_text = ' '.join(all_texts).strip()
                return combined_text
            except Exception as e:
                return self._unhashable_file_exception(file_info, e)

    def _extract_image_hash(self, file_info : FileInfo):
        """Extract image hash from a presentation, id it is possible."""
        self._handle_unsupported_ppt_exception(file_info)
        with suppress_stderr():
            try:
                phash_val = 0
                with zipfile.ZipFile(file_info.get_path(), 'r') as z:
                    img_files = [f for f in z.namelist() if f.startswith('ppt/media/')]
                    for _, fname in enumerate(img_files[:PDF_IMAGE_MAX_PAGES]):
                        with z.open(fname) as img_file:
                            img = Image.open(io.BytesIO(img_file.read()))
                            phash_val ^= self._get_phash_value(img)
                return phash_val
            except Exception as e:
                return self._unhashable_file_exception(file_info, e)

    def _unhashable_file_exception(self, file_info : FileInfo, error : Exception) -> int:
        """Handle fail to hash file exception."""
        self.logger.add_to_corrupted(file_info, error)
        file_info.set_auto_removability(False)
        return None

    def _handle_unsupported_ppt_exception(self, file_info : FileInfo) -> bool:
        """Return true, if specified file is unsupported type .ppt."""
        # skip .ppt presentations
        if file_info.get_suffix() == '.ppt':
            self.logger.add_to_corrupted(file_info, 'Unsupport presentation type .ppt')
            return True
        return False

            # combined_hash = int(hashlib.sha256(f'{texts_hash}{images_hash}'.encode()).hexdigest(), 16)
            # if combined_hash:
            #     return combined_hash
            # self.logger.add_to_corrupted(file_info, 'Failed to hash a file')
            # return None
