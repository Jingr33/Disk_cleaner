from pptx import Presentation
import hashlib
import zipfile
import io
from PIL import Image

from file_data.file_info import FileInfo
from hashers.text_hasher import TextHasher
from hashers.image_hasher import ImageHasher
from decorators import suppress_stderr
from config import PDF_IMAGE_MAX_PAGES

class PresentationHasher(TextHasher, ImageHasher):
    def __init__(self, sorter, logger):
        super().__init__(sorter, logger)

    def extract_hash(self, file_info : FileInfo) -> int:
        """Extract ppt file (presentation), return content as a string."""
        # skip .ppt presentations
        if file_info.get_suffix() == '.ppt':
            self.logger.add_to_corrupted(file_info, 'Unsupport presentation type .ppt')
            return None
        with suppress_stderr():
            texts_hash = self.extract_hash_from_text(file_info)
            images_hash = self.extract_hash_from_images(file_info)
            combined_hash = int(hashlib.sha256(f'{texts_hash}{images_hash}'.encode()).hexdigest(), 16)
            if combined_hash:
                return combined_hash
            self.logger.add_to_corrupted(file_info, 'Failed to hash a file')
            return None

    def extract_hash_from_text(self, file_info : FileInfo) -> int:
        """Extract hash from text in the presentation."""
        prs = Presentation(file_info.get_path())
        try:
            all_texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text:
                        all_texts.append(shape.text)
            combined_text = ' '.join(all_texts).strip()
            return super().extract_hash_from_text(combined_text)
        except Exception as e:
            self.logger.add_to_corrupted(file_info, e)
            file_info.set_auto_removability(False)
            return 0

    def extract_hash_from_images(self, file_info : FileInfo) -> int:
        """Extract hash from image database of presentation."""
        try:
            phash_val = 0
            with zipfile.ZipFile(file_info.get_path(), 'r') as z:
                img_files = [f for f in z.namelist() if f.startswith('ppt/media/')]
                for _, fname in enumerate(img_files[:PDF_IMAGE_MAX_PAGES]):
                    with z.open(fname) as img_file:
                        img = Image.open(io.BytesIO(img_file.read()))
                        phash_val ^= self.get_phash_value(img)
            return phash_val
        except Exception as e:
            self.logger.add_to_corrupted(file_info, e)
            file_info.set_auto_removability(False)
            return 0