import os
from pathlib import Path
from PIL import Image
import docx
import PyPDF2
from PyPDF2.errors import PdfReadError
import pandas as pd
from difflib import SequenceMatcher
import imagehash
import win32com.client
import pptx
import hashlib
from bs4 import BeautifulSoup
from zipfile import is_zipfile
from concurrent.futures import ThreadPoolExecutor, TimeoutError

# development
from icecream import ic

ROOT_FOLDER = r"D:\fyzio_disk\unzip"
MIN_SIMILARITY = 0.8 # above this similarity, files will be evaluated as similar
AUTO_REMOVE_SIMILARITY = 1.0 # above this simililarity, the second file will be remove automaticly


FILE_TYPES = ["text", "docx", "doc", "pdf", "spreadsheet", "pptx", "image", "htm", 'other']

# HASHES
def extract_text_file(path: Path) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

def is_probably_valid_docx(path: Path) -> bool:
    return path.suffix.lower() == ".docx" and is_zipfile(path)
    
def extract_doc_file(path: Path) -> str:
    try:
        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False
        doc = word.Documents.Open(str(path))
        try:
            text = doc.Content.Text
        finally:
            doc.Close(False)
            word.Quit()

        cleaned_text = text.strip().replace("\r", "").replace("\n", "")
        return hashlib.sha256(cleaned_text.encode('utf-8')).hexdigest()
    
    except Exception as e:
        print(f"Chyba při čtení DOC ({path.name}): {e}")
        return ""

def extract_docx_file(path: Path) -> str:
    try:
        if is_probably_valid_docx(path):
            doc = docx.Document(path)
            return "\n".join([para.text for para in doc.paragraphs])
        else:
            print(f"Soubor {path} není validní .docx, zkusíme jako .doc...")
            return None
    except Exception as e:
        print(f"Selhalo zpracování souboru {path}: {e}")
        return None

def _read_pdf(path: Path) -> str:
    text = ""
    with open(path, "rb") as f:
        reader = PyPDF2.PdfReader(f)
        for page in reader.pages:
            text += page.extract_text() or ""
    return text

def extract_pdf_file(path: Path) -> str:
    try:
        with ThreadPoolExecutor(max_workers=1) as executor:
            future = executor.submit(_read_pdf, path)
            return future.result(timeout=4)    
    except TimeoutError:
        print(f"Timeout při čtení PDF: {path}")
        return None
    except PdfReadError as e:
        print(f"Chyba při čtení PDF ({path}): {e}")
        return None
    except Exception as e:
        print(f"Jiná chyba při PDF ({path}): {e}")
        return None

def extract_spreadsheet_file(path: Path) -> str:
    try:
        if path.suffix.lower() == ".csv":
            df = pd.read_csv(path)
        else:
            df = pd.read_excel(path, engine='openpyxl')
        return df.to_csv(index=False)
    except:
        return None
    
def extract_ppt_file(path: Path) -> str:
    try:
        prs = pptx.Presentation(path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)
    except Exception as e:
        return f"ERROR: {e}"

def extract_htm_file(path: Path) -> str:
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            soup = BeautifulSoup(f, "html.parser")
        return soup.get_text()
    except Exception as e:
        return f"ERROR: {e}"

def get_image_phash(path: Path):
    try:
        hash = imagehash.average_hash(Image.open(path))
        return str(hash)
    except Exception as e:
        print(f"Chyba při zpracování {path}: {e}")
        return None
 
def count_hashes(sorted_files : dict) -> dict:
    print("Counting file hashes...")
    hashed_files = {}
    for type in sorted_files.keys():
        hashed_files[type] = []

    hash_functions = {
        "text" : extract_text_file,
        "docx" : extract_docx_file,
        "doc" : extract_doc_file,
        "pdf" : extract_pdf_file,
        "spreadsheet" : extract_spreadsheet_file,
        "pptx" : extract_ppt_file,
        "htm" : extract_htm_file,
        "image" : get_image_phash,
    }

    counter = 0
    for file_type in sorted_files.keys():
        for file_path in sorted_files[file_type]:
            counter += 1
            print(counter)
            if not file_path.exists() and not file_path.is_files():
                continue
            file_hash = hash_functions[file_type](file_path)
            hashed_files[file_type].append((file_path, file_hash))
    print("Hashes completed...")
    return hashed_files

# COMPARSION
def similarity_score(hash1: str, hash2: str) -> float:
    if hash1 is None or hash2 is None:
        return float(0)
    return SequenceMatcher(None, hash1, hash2).ratio()

# FILE COLLECTION  
def explore_folder(folder_path : str, files : list) -> list:
    """Find all files inside of the root folder."""
    for item in os.listdir(folder_path):
        item_path = Path(os.path.join(folder_path, item))
        if os.path.isdir(item_path):
            files = explore_folder(item_path, files)
        else:
            files.append(item_path)
    return files

def remove_wave_starters(file_paths : list) -> list:
    print("\n\n\n")
    for i in range(len(file_paths) -1, -1, -1):
        path = Path(file_paths[i])
        if '~' in path.name:
            path.unlink()
            del file_paths[i]
            print(f"\nSoubor {path.name} byl odstraněn.")
    print("\n\n\n")

def sort_by_extension(unsorted_files : list) -> dict:
    """Sort files into groups by file type."""
    sorted_files = {}
    for type in FILE_TYPES:
        sorted_files[type] = []

    for file_path in unsorted_files:
        file_extension = file_path.suffix
        if '~' in file_path.name:
            continue

        if file_extension in ['.txt', '.md']:
            sorted_files['text'].append(file_path)
        elif file_extension == '.docx':
            sorted_files['docx'].append(file_path)
        elif file_extension in ['.doc', '.docm']:
            sorted_files['doc'].append(file_path)
        elif file_extension == '.pdf':
            sorted_files['pdf'].append(file_path)
        elif file_extension in ['.ppt', '.pptx']:
            sorted_files['pptx'].append(file_path)
        elif file_extension in ['.csv', '.xlsx', '.xls']:
            sorted_files["spreadsheet"].append(file_path)
        elif file_extension.lower() in ['.png', '.jpg', '.jpeg', '.bmp', '.gif', '.tif', '.heic']:
            sorted_files['image'].append(file_path)
        elif file_extension == '.htm':
            sorted_files['htm'].append(file_path)
        else:
            sorted_files['other'].append(file_path)
    return sorted_files

def print_file_counts(sorted_files : dict) -> None:
    print("\nPOČTY TYPŮ SOUBORŮ:")
    for type in sorted_files.keys():
        type_count =  len(sorted_files[type])
        print(f"{type.upper()}: {type_count}")

def select_entered_file_types(sorted_files : dict) -> None:
    user_input = input("JAKÝ TYP SOUBOR SE MÁ VYTŘÍDIT? (text, docx, doc, pptx, pdf, spreadsheet, image, html)\n")
    keys_to_delete = []

    for file_type in list(sorted_files.keys()):
        if file_type not in user_input.lower() or not sorted_files[file_type]:
            keys_to_delete.append(file_type)

    for key in keys_to_delete:
        del sorted_files[key]

    return sorted_files

#REMOVING
def manage_removing(hashed_files : dict) -> None:
    for ftype in  hashed_files.keys():
        (file1, hash1) = hashed_files[ftype][-1]
        for idx_file2 in range(len(hashed_files[ftype]) - 2, -1, -1):

            if len(hashed_files[ftype]) < 2:
                break

            (file2, hash2) = hashed_files[ftype][idx_file2]
            sim_score = similarity_score(hash1, hash2)
            print(sim_score)

            if sim_score < MIN_SIMILARITY:
                continue
            if not file1.exists() or not file1.is_file():
                del hashed_files[ftype][-1]
                continue
            if not file2.exists() or not file2.is_file():
                del hashed_files[ftype][idx_file2]
                continue

            if sim_score >= AUTO_REMOVE_SIMILARITY:
                hashed_files[ftype][idx_file2][0].unlink()
                print(f"\nSoubor {file2} byl odstraněn.")
            else:
                hashed_files[ftype] = ask_for_remove(hashed_files[ftype], file1, file2, sim_score, idx_file2)

def ask_for_remove(files : list, fpath1 : Path, fpath2 : Path, sim_score, idx_file2 : int) -> list:
    print(f"\nPodobnostní skóre je {round(sim_score * 100, 2)} % mezi {fpath1.name} a {fpath2.name} soubory.")
    os.startfile(fpath1)
    os.startfile(fpath2)
    user_input = input(f"Chcete odebrat {fpath2.name}? (Y/n)")
    if not user_input == "Y":
        print(f"{fpath2} byl zachován.")
        return files
    
    while is_file_locked(fpath1) or is_file_locked(fpath2):
        input("One of the files is still open. Close it and press enter for continue.")
    fpath2.unlink()
    del files[idx_file2]
    print(f"{fpath2} byl odstraněn.")
    return files

def is_file_locked(path: Path) -> bool:
    try:
        with open(path, "a"):
            return False
    except IOError:
        return True
    
def find_same_names(files : list) -> list:
    counter = 0
    duplicities = []
    for i1 in range(len(files)):
        for i2 in range(i1 + 1, len(files)):
            path1 = Path(files[i1])
            path2 = Path(files[i2])

            if path1.name != path2.name:
                continue

            duplicities.append((path1, path2))
            counter += 1
    print(counter)
    return duplicities

def remove_duplicate_name_files(name_duplicities : list[tuple]) -> None:
    remove_all_automaticly = False
    while len(name_duplicities):
        (path1, path2) = name_duplicities[0]

        if not files_exists((name_duplicities[0])):
            del name_duplicities[0]
            continue

        print(f"\nVícečetný název:\n{path1.name}\n{path2.name}")
        if not remove_all_automaticly:
            user_input = input("Do you want to remove second file ?\n(Y - remove, n - keep file, All - remove all duplicities)")

        if user_input == "All" or remove_all_automaticly:
                name_duplicities = remove_duplicate_name_file(name_duplicities, 0)
                remove_all_automaticly = True
        elif user_input == "Y":
            name_duplicities = remove_duplicate_name_file(name_duplicities, 0)
        else:
            del name_duplicities[0]

    print("\nDuplicate names are completed.")

def remove_duplicate_name_file(duplicities : list[tuple], index : int) -> list[tuple]:
    duplicity = duplicities.pop(index)
    duplicity[1].unlink()
    print(f"file {duplicity[1]} was removed.")
    return duplicities

def files_exists(paths : tuple) -> bool:
    for path in paths:
        if not path.exists(): 
            return False
    return True

all_files = explore_folder(ROOT_FOLDER, [])
# all_files = remove_wave_starters(all_files)
# sorted_files = sort_by_extension(all_files)
# print_file_counts(sorted_files)
# sorted_files = select_entered_file_types(sorted_files)
# hashed_files = count_hashes(sorted_files)
# # removing
# manage_removing(hashed_files)

duplicities = find_same_names(all_files)
remove_duplicate_name_files(duplicities)